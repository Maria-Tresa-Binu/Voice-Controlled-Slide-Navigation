import streamlit as sl
from spire.presentation.common import*
import spire.presentation as sp
from pptx import Presentation
#import SpeechRecognition

def read_ppt(file):
    ppt = Presentation(file)
    for slide in ppt.slides:
        text=[]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        if text[0]!='<#>':
            sl.markdown("\n".join(text))
            sl.markdown("---")        

    #return text

def main():
    
    sl.title("Upload your file here")
    sl.markdown("---")
    presentation=sl.file_uploader("please upload your ppt here",type=["pptx"])
    if presentation is not None:
        read_ppt(presentation)
        #ppt_to_img(presentation)
        
#if __name__ == "__main__":
main()
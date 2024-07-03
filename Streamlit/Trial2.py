import streamlit as st

from pptx import Presentation
from io import BytesIO
def extract_text(ppt):
    prs=Presentation(ppt)
    text_content=[]
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text_content.append(shape.text)
    return "\n".join(text_content)
def main():
    st.title("PPT Text Extractor")
    st.write("upload your ppt")
    file=st.file_uploader('Choose a ppt file',type='pptx')

    if file is not None:
        text=extract_text(file)
        st.header("Extracted Text")
        st.text_area("Text Content",text,height=400)

if __name__=='__main__':
    main()